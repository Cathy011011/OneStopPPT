from pathlib import Path
import shutil

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_VERTICAL_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(r"F:\Codex\PPT自动化")
INTERVIEW_ROOT = ROOT / "访谈"
REF_DIR = INTERVIEW_ROOT / "01_参考模板"
INPUT_DIR = INTERVIEW_ROOT / "02_访谈资料"
OUTPUT_DIR = INTERVIEW_ROOT / "03_交付成品"
OUT = OUTPUT_DIR / "座谈会笔录2-7月13日_参考模板可编辑版.pptx"
PREVIEW_DIR = OUTPUT_DIR / "0713-reference-editable-preview"
FONT = "Microsoft YaHei"
FONT_NUM = "Arial"
SOURCE = "资料来源：7月13日用户座谈会纪要；按归纳总结提示词重组。"

C = {
    "white": RGBColor(0xFF, 0xFF, 0xFF),
    "bg": RGBColor(0xFB, 0xFB, 0xFB),
    "dark": RGBColor(0x18, 0x18, 0x18),
    "charcoal": RGBColor(0x3E, 0x3E, 0x38),
    "gray": RGBColor(0x71, 0x71, 0x71),
    "light": RGBColor(0xE9, 0xEC, 0xED),
    "line": RGBColor(0xC9, 0xCE, 0xD1),
    "orange": RGBColor(0xF0, 0x61, 0x24),
    "teal": RGBColor(0x0C, 0xB6, 0xA6),
    "blue": RGBColor(0x22, 0x74, 0xAC),
    "green": RGBColor(0x4B, 0x8D, 0x56),
    "red": RGBColor(0xC0, 0x00, 0x00),
    "soft_orange": RGBColor(0xFC, 0xF0, 0xEA),
    "soft_teal": RGBColor(0xEB, 0xF7, 0xF4),
    "soft_blue": RGBColor(0xEE, 0xF4, 0xFA),
    "soft_gray": RGBColor(0xF3, 0xF5, 0xF6),
}


def set_bg(slide, color):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_rect(slide, left, top, width, height, fill, line=None, line_width=0.8, transparency=0):
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    shape.fill.transparency = transparency
    shape.line.color.rgb = line or fill
    shape.line.width = Pt(line_width)
    return shape


def add_line(slide, left, top, width, color, height=0.02):
    return add_rect(slide, left, top, width, height, color, color, 0)


def add_text(
    slide,
    left,
    top,
    width,
    height,
    text,
    font_size=16,
    color=None,
    bold=False,
    align=PP_ALIGN.LEFT,
    valign=MSO_VERTICAL_ANCHOR.TOP,
    font_name=FONT,
    margin=2,
):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    tf.vertical_anchor = valign
    tf.margin_left = Pt(margin)
    tf.margin_right = Pt(margin)
    tf.margin_top = Pt(margin)
    tf.margin_bottom = Pt(margin)
    p = tf.paragraphs[0]
    p.alignment = align
    p.space_before = 0
    p.space_after = 0
    run = p.add_run()
    run.text = text
    run.font.name = font_name
    run.font.size = Pt(font_size)
    run.font.bold = bold
    run.font.color.rgb = color or C["charcoal"]
    return box


def add_bullets(slide, left, top, width, items, font_size=16, color=None, bullet_color=None, gap=0.34):
    y = top
    for item in items:
        dot = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(y + 0.09), Inches(0.10), Inches(0.10)
        )
        dot.fill.solid()
        dot.fill.fore_color.rgb = bullet_color or C["orange"]
        dot.line.color.rgb = bullet_color or C["orange"]
        add_text(slide, left + 0.16, y, width - 0.16, 0.28, item, font_size, color=color or C["charcoal"])
        y += gap


def add_table(slide, left, top, width, height, headers, rows, col_widths, font_size=12, header_fill=None):
    shape = slide.shapes.add_table(len(rows) + 1, len(headers), Inches(left), Inches(top), Inches(width), Inches(height))
    table = shape.table
    total = sum(col_widths)
    for idx, w in enumerate(col_widths):
        table.columns[idx].width = Inches(width * (w / total))

    for c, header in enumerate(headers):
        cell = table.cell(0, c)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = header_fill or C["charcoal"]
        cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
        for p in cell.text_frame.paragraphs:
            p.alignment = PP_ALIGN.CENTER
            for run in p.runs:
                run.font.name = FONT
                run.font.size = Pt(font_size)
                run.font.bold = True
                run.font.color.rgb = C["white"]

    for r, row in enumerate(rows, start=1):
        for c, value in enumerate(row):
            cell = table.cell(r, c)
            cell.text = value
            cell.fill.solid()
            cell.fill.fore_color.rgb = C["white"] if r % 2 else C["soft_gray"]
            cell.vertical_anchor = MSO_VERTICAL_ANCHOR.MIDDLE
            cell.text_frame.word_wrap = True
            for p in cell.text_frame.paragraphs:
                p.alignment = PP_ALIGN.LEFT
                for run in p.runs:
                    run.font.name = FONT
                    run.font.size = Pt(font_size)
                    run.font.bold = c == 0
                    run.font.color.rgb = C["charcoal"]
    return table


def page_chrome(slide, page_no, title, kicker=None):
    set_bg(slide, C["bg"])
    add_line(slide, 0.56, 0.46, 12.2, C["charcoal"], 0.02)
    add_rect(slide, 0.56, 0.68, 0.08, 0.54, C["orange"], C["orange"], 0)
    if kicker:
        add_text(slide, 0.76, 0.68, 2.4, 0.18, kicker, 9.5, color=C["gray"], bold=True, font_name=FONT_NUM)
    add_text(slide, 0.76, 0.90, 11.2, 0.34, title, 24, color=C["charcoal"], bold=True)
    add_line(slide, 0.76, 1.34, 1.25, C["orange"], 0.03)
    add_line(slide, 0.56, 6.92, 12.2, C["line"], 0.01)
    add_text(slide, 0.56, 6.96, 9.0, 0.14, SOURCE, 9, color=C["gray"])
    add_text(slide, 12.10, 6.96, 0.45, 0.14, f"{page_no:02d}", 9, color=C["gray"], align=PP_ALIGN.RIGHT, font_name=FONT_NUM)


def section_slide(slide, page_no, no, title, desc, items):
    set_bg(slide, C["dark"])
    add_rect(slide, 0.86, 0.96, 0.10, 5.18, C["orange"], C["orange"], 0)
    add_text(slide, 1.18, 1.22, 1.10, 0.70, no, 40, color=C["white"], bold=True, font_name=FONT_NUM)
    add_text(slide, 2.04, 1.46, 8.80, 0.48, title, 28, color=C["white"], bold=True)
    add_text(slide, 2.06, 2.18, 8.20, 0.44, desc, 15, color=RGBColor(0xD6, 0xDC, 0xE1))
    y = 3.18
    for idx, item in enumerate(items, start=1):
        add_rect(slide, 2.06, y + 0.03, 0.34, 0.28, C["orange"], C["orange"], 0)
        add_text(slide, 2.12, y + 0.02, 0.22, 0.16, f"{idx}", 9.5, color=C["white"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_NUM)
        add_text(slide, 2.48, y, 7.60, 0.22, item, 13.5, color=RGBColor(0xD6, 0xDC, 0xE1))
        y += 0.42
    add_line(slide, 0.84, 6.74, 11.50, RGBColor(0x77, 0x83, 0x86), 0.01)
    add_text(slide, 12.02, 6.76, 0.28, 0.14, f"{page_no:02d}", 9, color=C["white"], align=PP_ALIGN.RIGHT, font_name=FONT_NUM)


def metric_card(slide, left, top, width, height, title, note, accent, fill=None):
    add_rect(slide, left, top, width, height, fill or C["white"], C["line"], 0.8)
    add_rect(slide, left, top, width, 0.14, accent, accent, 0)
    add_text(slide, left + 0.14, top + 0.22, width - 0.28, 0.22, title, 15.5, color=C["charcoal"], bold=True)
    add_text(slide, left + 0.14, top + 0.50, width - 0.28, height - 0.62, note, 11.5, color=C["gray"])


def matrix_page(slide, page_no, title, top_labels, side_labels, body, colors):
    page_chrome(slide, page_no, title, "MATRIX")
    x0 = 1.00
    y0 = 1.82
    left_w = 1.95
    col_w = 2.95
    head_h = 0.56
    row_h = 1.18
    for i, label in enumerate(top_labels):
        add_rect(slide, x0 + left_w + col_w * i, y0, col_w - 0.02, head_h, colors[i], colors[i], 0)
        add_text(slide, x0 + left_w + col_w * i + 0.06, y0 + 0.12, col_w - 0.14, 0.18, label, 13, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
    for r, side in enumerate(side_labels):
        add_rect(slide, x0, y0 + head_h + row_h * r, left_w, row_h - 0.02, C["soft_gray"], C["white"], 0)
        add_text(slide, x0 + 0.10, y0 + head_h + row_h * r + 0.34, left_w - 0.18, 0.30, side, 14, color=C["charcoal"], bold=True, align=PP_ALIGN.CENTER)
        for c, text in enumerate(body[r]):
            fill = C["white"] if (r + c) % 2 == 0 else C["soft_gray"]
            add_rect(slide, x0 + left_w + col_w * c, y0 + head_h + row_h * r, col_w - 0.02, row_h - 0.02, fill, C["white"], 0)
            add_text(slide, x0 + left_w + col_w * c + 0.12, y0 + head_h + row_h * r + 0.12, col_w - 0.24, row_h - 0.18, text, 12.5, color=C["charcoal"])


def bar_compare(slide, left, top, width, label, level, color):
    add_text(slide, left, top, 1.0, 0.20, label, 13, color=C["charcoal"], bold=True)
    add_rect(slide, left + 1.10, top + 0.05, width, 0.12, C["light"], C["light"], 0)
    add_rect(slide, left + 1.10, top + 0.05, width * level, 0.12, color, color, 0)
    tip = "高频" if level > 0.72 else "中频" if level > 0.52 else "待验证"
    add_text(slide, left + 1.10 + width + 0.12, top - 0.01, 0.70, 0.18, tip, 11, color=C["gray"])


def export_preview(path):
    import win32com.client

    if PREVIEW_DIR.exists():
        shutil.rmtree(PREVIEW_DIR)
    PREVIEW_DIR.mkdir(parents=True, exist_ok=True)
    app = win32com.client.DispatchEx("PowerPoint.Application")
    app.Visible = 1
    deck = app.Presentations.Open(str(path), WithWindow=False)
    deck.SaveAs(str(PREVIEW_DIR), 18)
    deck.Close()
    app.Quit()


def build():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    blank = prs.slide_layouts[6]

    slide = prs.slides.add_slide(blank)
    set_bg(slide, C["bg"])
    add_rect(slide, 0.00, 1.10, 13.333, 2.28, RGBColor(0x55, 0x66, 0x6B), RGBColor(0x55, 0x66, 0x6B), 0)
    add_rect(slide, 0.00, 1.10, 13.333, 2.28, RGBColor(0x18, 0x18, 0x18), RGBColor(0x18, 0x18, 0x18), 0, 0.30)
    for i in range(10):
        add_rect(slide, i * 1.4 - 0.8, 1.14, 0.03, 2.18, RGBColor(0xE3, 0xE8, 0xEA), RGBColor(0xE3, 0xE8, 0xEA), 0, 0.65)
    add_rect(slide, 0.92, 0.84, 4.18, 0.10, RGBColor(0x7B, 0x88, 0x8D), RGBColor(0x7B, 0x88, 0x8D), 0)
    add_text(slide, 2.18, 2.06, 9.0, 0.38, "7月13日用户座谈会研究总结", 28, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 2.24, 2.52, 8.9, 0.24, "基于访谈纪要重组为研究汇报口径；保持原生可编辑对象", 12, color=RGBColor(0xE4, 0xE8, 0xEA), align=PP_ALIGN.CENTER)
    metric_card(slide, 1.05, 4.22, 3.30, 1.28, "会议形式", "用户座谈会，讨论从汽车态度到购车决策与服务体验。", C["teal"])
    metric_card(slide, 5.00, 4.22, 3.30, 1.28, "样本结构", "6位新能源用户，覆盖纯电、插混、增程及多场景用车。", C["orange"])
    metric_card(slide, 8.95, 4.22, 3.30, 1.28, "输出口径", "用于内部产品、设计、服务和传播协同的研究总结。", C["blue"])
    add_text(slide, 0.96, 6.78, 9.4, 0.16, SOURCE, 9, color=C["gray"])

    slide = prs.slides.add_slide(blank)
    set_bg(slide, C["bg"])
    add_rect(slide, 0.00, 0.00, 13.333, 4.44, RGBColor(0x38, 0x82, 0xA8), RGBColor(0x38, 0x82, 0xA8), 0)
    for i in range(11):
        add_rect(slide, i * 1.18 - 0.5, 0.00, 0.04, 4.44, RGBColor(0xFF, 0xFF, 0xFF), RGBColor(0xFF, 0xFF, 0xFF), 0, 0.78)
    add_rect(slide, 5.34, 0.22, 2.32, 1.42, RGBColor(0xE5, 0x83, 0x51), RGBColor(0xE5, 0x83, 0x51), 0)
    add_text(slide, 5.34, 0.62, 2.32, 0.26, "目录", 20, color=C["white"], bold=True, align=PP_ALIGN.CENTER)
    add_text(slide, 5.34, 0.98, 2.32, 0.20, "Contents", 10.5, color=C["white"], align=PP_ALIGN.CENTER, font_name=FONT_NUM)
    toc_cols = [
        ("1", "会议与证据框架", ["会议概况", "关键结论", "决策链路"]),
        ("2", "用户与竞品判断", ["需求矩阵", "品牌门槛", "车型短板"]),
        ("3", "建议与落地方向", ["机会点", "服务数字化", "行动路径"]),
    ]
    x_positions = [0.92, 4.52, 8.12]
    for x, (num, title, items) in zip(x_positions, toc_cols):
        add_rect(slide, x, 3.12, 3.16, 2.32, C["white"], C["white"], 0, 0.06)
        circ = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(x + 1.28), Inches(2.52), Inches(0.58), Inches(0.58))
        circ.fill.solid()
        circ.fill.fore_color.rgb = C["teal"]
        circ.line.color.rgb = C["white"]
        add_text(slide, x + 1.41, 2.66, 0.30, 0.16, num, 10, color=C["white"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_NUM)
        add_text(slide, x + 0.20, 3.32, 2.76, 0.22, title, 13.5, color=C["charcoal"], bold=True, align=PP_ALIGN.CENTER)
        y = 3.76
        for item in items:
            add_text(slide, x + 0.22, y, 2.72, 0.18, item, 10.5, color=C["charcoal"])
            y += 0.22
    add_text(slide, 12.36, 6.95, 0.30, 0.14, "02", 9, color=C["gray"], align=PP_ALIGN.RIGHT, font_name=FONT_NUM)

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 3, "项目概况：这次讨论覆盖了新能源用户从认知到决策的完整链路", "PROJECT SUMMARY")
    add_table(
        slide,
        0.86,
        1.70,
        4.15,
        2.80,
        ["维度", "内容"],
        [
            ["会议形式", "用户座谈会"],
            ["样本数量", "6位受访者"],
            ["样本结构", "纯电 / 插混 / 增程并存"],
            ["核心主题", "汽车态度、购车逻辑、品牌与车型判断"],
            ["输出方式", "研究总结与行动建议"],
        ],
        [1.1, 2.2],
        11,
        C["charcoal"],
    )
    add_table(
        slide,
        5.28,
        1.70,
        6.92,
        2.80,
        ["主线", "用户关注", "对汇报意味着什么"],
        [
            ["汽车态度", "安全、智能、个性与生活方式", "产品不能只讲交通工具属性"],
            ["新能源驱动", "省成本 + 动力快 + 更安静", "要表达成套收益，而非单点参数"],
            ["决策框架", "预算、品牌、续航、空间、动力、服务", "销售现场更像验证，而不是重塑"],
            ["产品判断", "基础认可存在，但购买理由不够强", "品牌与产品体验需要同时补齐"],
        ],
        [1.0, 2.1, 3.4],
        11,
        C["orange"],
    )
    metric_card(slide, 0.86, 4.88, 2.75, 1.20, "结论 1", "用户已把汽车视作生活方式和情绪价值载体。", C["teal"], C["soft_teal"])
    metric_card(slide, 3.84, 4.88, 2.75, 1.20, "结论 2", "新能源吸引力来自真实成本与体验升级。", C["orange"], C["soft_orange"])
    metric_card(slide, 6.82, 4.88, 2.75, 1.20, "结论 3", "传统品牌仍要重新解释新能源价值。", C["blue"], C["soft_blue"])
    metric_card(slide, 9.80, 4.88, 2.40, 1.20, "结论 4", "服务体验已经进入车型竞争力。", C["red"], C["soft_gray"])

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 4, "核心发现：用户希望汽车同时满足安全感、智能化和生活方式表达", "KEY FINDINGS")
    add_rect(slide, 0.90, 1.72, 3.22, 4.80, C["soft_blue"], C["line"], 0.8)
    add_text(slide, 1.12, 2.00, 2.78, 0.74, "汽车不再只是代步工具，而是家庭安全、个人审美和日常情绪价值的综合载体。", 20, color=C["charcoal"], bold=True)
    add_bullets(
        slide,
        1.16,
        3.34,
        2.76,
        [
            "安全是有家庭场景用户的前置门槛",
            "智能化已成为当代汽车的基础能力",
            "有人希望通过改装和设计表达自己的偏好",
            "“只是交通工具”已无法解释真实购买动机",
        ],
        12.5,
        C["charcoal"],
        C["orange"],
        0.46,
    )
    add_text(slide, 4.64, 1.84, 3.20, 0.22, "关注议题层级", 16, color=C["charcoal"], bold=True)
    bar_compare(slide, 4.64, 2.28, 2.85, "安全", 0.86, C["teal"])
    bar_compare(slide, 4.64, 2.72, 2.85, "智能", 0.80, C["orange"])
    bar_compare(slide, 4.64, 3.16, 2.85, "空间", 0.72, C["blue"])
    bar_compare(slide, 4.64, 3.60, 2.85, "审美", 0.67, C["teal"])
    bar_compare(slide, 4.64, 4.04, 2.85, "舒适", 0.63, C["orange"])
    bar_compare(slide, 4.64, 4.48, 2.85, "改装表达", 0.44, C["blue"])
    add_table(
        slide,
        7.90,
        1.72,
        4.30,
        4.80,
        ["观察", "对产品意味着什么"],
        [
            ["“买车就得买安全的”", "安全不是加分项，而是天然门槛"],
            ["“希望我的车子很智能”", "智能体验要能被用户直接感知"],
            ["“喜欢把车改成自己喜欢的样子”", "设计与个性表达会影响长期好感"],
            ["家庭/通勤/旅行场景并存", "产品要覆盖多场景，而非单一定位"],
        ],
        [1.7, 2.6],
        10.5,
        C["charcoal"],
    )

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 5, "决策链路：预算只是起点，真正进入名单还要同时过品牌、空间和体验三道门", "DECISION PROCESS")
    steps = [
        ("01", "预算范围", "先圈定价格带"),
        ("02", "品牌门槛", "口碑、成熟度、技术来历"),
        ("03", "外观内饰", "风格、质感与第一印象"),
        ("04", "空间续航", "家庭与长途场景能否覆盖"),
        ("05", "试驾服务", "动力、讲解和政策透明度"),
    ]
    x = 0.90
    for idx, (no, title, note) in enumerate(steps):
        add_rect(slide, x, 2.08, 2.26, 1.78, C["white"], C["line"], 0.8)
        add_rect(slide, x, 2.08, 2.26, 0.14, C["orange"] if idx % 2 == 0 else C["teal"], C["orange"] if idx % 2 == 0 else C["teal"], 0)
        add_text(slide, x + 0.12, 2.28, 0.36, 0.14, no, 9.5, color=C["gray"], bold=True, font_name=FONT_NUM)
        add_text(slide, x + 0.12, 2.54, 1.96, 0.18, title, 15, color=C["charcoal"], bold=True, align=PP_ALIGN.CENTER)
        add_text(slide, x + 0.12, 3.02, 1.96, 0.36, note, 11.5, color=C["gray"], align=PP_ALIGN.CENTER)
        if idx < len(steps) - 1:
            add_line(slide, x + 2.26, 2.96, 0.42, C["gray"], 0.03)
        x += 2.48
    add_table(
        slide,
        0.90,
        4.38,
        11.32,
        2.12,
        ["环节", "用户判断方式", "对应提醒"],
        [
            ["预算/品牌", "看价格带，也看谁用得多、是否成熟可靠", "品牌负责让用户愿意进一步了解"],
            ["产品体验", "比外观、内饰、空间、续航和动力", "进入名单后，体验差异才决定是否留下"],
            ["现场验证", "试驾、服务、金融和政策透明度", "销售不是认知起点，而是最终验证阶段"],
        ],
        [1.1, 2.9, 3.0],
        11,
        C["orange"],
    )

    slide = prs.slides.add_slide(blank)
    section_slide(
        slide,
        6,
        "02",
        "用户与竞品判断",
        "从用户需求矩阵到品牌进入门槛，再到车型被排除的原因。",
        ["多场景需求如何同时影响购买", "品牌可信度如何转化为新能源心智", "为什么基础认可没有形成购买闭环"],
    )

    slide = prs.slides.add_slide(blank)
    matrix_page(
        slide,
        7,
        "需求矩阵：用户期待一台能同时覆盖家庭、通勤和旅行的生活载体",
        ["安全与安心", "智能与便利", "审美与情绪"],
        ["家庭出行", "城市通勤", "长途旅行"],
        [
            ["安全是家庭场景的前置门槛", "辅助体验与控车便利要易感知", "空间与舒适度提升陪伴感"],
            ["可靠、耐用、成本可预期", "安静、灵活、数字化使用顺畅", "车辆也承载个人风格表达"],
            ["真实续航与补能安心", "导航、车机与远程功能更重要", "久坐舒适和车内氛围影响旅行感受"],
        ],
        [C["orange"], C["teal"], C["blue"]],
    )

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 8, "品牌能帮助车型进入第一轮名单，但新能源时代仍需重新建立独立说服力", "BRAND JUDGMENT")
    add_rect(slide, 0.92, 1.74, 5.48, 4.88, C["white"], C["line"], 0.8)
    add_text(slide, 1.16, 2.02, 2.40, 0.20, "品牌进入门槛", 16, color=C["charcoal"], bold=True)
    add_bullets(
        slide,
        1.16,
        2.48,
        4.90,
        [
            "用户用“谁用得多、是否成熟可靠、核心技术是否可信”判断品牌",
            "朋友推荐、市场保有量和技术口碑会共同构成可信度",
            "传统品牌的燃油车声量，不会自动迁移成新能源心智",
            "进入新能源赛道后，用户仍会重新追问：它是不是也足够强",
        ],
        12.5,
        C["charcoal"],
        C["teal"],
        0.46,
    )
    add_rect(slide, 6.72, 1.74, 5.48, 4.88, C["soft_gray"], C["line"], 0.8)
    add_text(slide, 6.96, 2.02, 2.60, 0.20, "代表性感知", 16, color=C["charcoal"], bold=True)
    metric_card(slide, 6.96, 2.38, 2.38, 1.14, "比亚迪", "口碑好、性价比高、电池技术可信", C["orange"], C["white"])
    metric_card(slide, 9.58, 2.38, 2.38, 1.14, "特斯拉", "安全性高、操控好、保值率较强", C["teal"], C["white"])
    metric_card(slide, 6.96, 3.74, 2.38, 1.14, "理想", "舒适、空间更大、车机体验更丰富", C["blue"], C["white"])
    metric_card(slide, 9.58, 3.74, 2.38, 1.14, "某合资品牌", "传统品牌可信，但新能源认知尚未完全建立", C["charcoal"], C["white"])
    add_text(slide, 6.98, 5.30, 4.90, 0.42, "结论：品牌负责建立“愿意去看”的门槛，车型体验才决定是否成交。", 13.5, color=C["charcoal"], bold=True)

    slide = prs.slides.add_slide(blank)
    matrix_page(
        slide,
        9,
        "车型被排除的核心，不是没有基础认可，而是没有形成足够强的新能源购买理由",
        ["基础认可", "主要短板", "汇报判断"],
        ["品牌", "产品", "服务"],
        [
            ["传统品牌可靠、价格大体可接受", "新能源心智不完整", "基础盘存在，但不能替代产品说服"],
            ["有空间和一定配置讨论空间", "续航、动力、车机和定位不够强", "没有形成区别于竞品的清晰理由"],
            ["用户愿意进一步了解", "讲参数、缺少温度与顾客视角", "服务体验已进入车型竞争力"],
        ],
        [C["orange"], C["teal"], C["blue"]],
    )
    add_text(slide, 1.00, 6.18, 11.00, 0.18, "注：具体设计部位、配置梯度及优先级排序仍需问卷进一步验证。", 10.5, color=C["gray"])

    slide = prs.slides.add_slide(blank)
    section_slide(
        slide,
        10,
        "03",
        "建议与落地方向",
        "围绕续航可信、动力爽感、设计科技感和顾客视角服务，形成后续动作。",
        ["产品与座舱如何回答用户的第一眼判断", "服务与数字化如何承接购车链路", "哪些问题需要进入下一轮问卷验证"],
    )

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 11, "机会点：把用户不满转化为更明确的产品切入点", "OPPORTUNITY")
    metric_card(slide, 0.96, 1.82, 2.65, 1.42, "真实续航与补能安心", "把“续航短/虚高”的不安感，转为可解释、可验证的使用信心。", C["orange"], C["soft_orange"])
    metric_card(slide, 3.82, 1.82, 2.65, 1.42, "新能源动力爽感", "强化起步、加速、超车时的直接响应，让试驾更有记忆点。", C["teal"], C["soft_teal"])
    metric_card(slide, 6.68, 1.82, 2.65, 1.42, "更年轻的内饰风格", "减少商务感和单调感，让座舱更清晰、更科技、更一致。", C["blue"], C["soft_blue"])
    metric_card(slide, 9.54, 1.82, 2.65, 1.42, "更有温度的服务", "让销售从“背参数”转成“讲场景价值”，减少体验流失。", C["charcoal"], C["soft_gray"])
    add_table(
        slide,
        0.96,
        3.64,
        11.24,
        2.84,
        ["问题来源", "用户感受", "产品/传播切入点"],
        [
            ["续航与补能", "担心不够实、不够安心", "用真实场景解释续航与补能边界"],
            ["动力感知", "“不太像新能源”", "试驾突出提速、静音和平顺的差异感"],
            ["座舱体验", "像商务车、车机比较单调", "统一设计、车机与氛围表达，建立第一眼记忆"],
            ["服务体验", "只有参数、不够有温度", "增加顾客场景代入、对比解释和试驾引导"],
        ],
        [1.4, 2.1, 3.3],
        11,
        C["charcoal"],
    )

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 12, "设计与配置建议：先补齐基础体验，再谈更细的配置梯度", "PRODUCT ACTION")
    add_rect(slide, 0.96, 1.80, 5.38, 4.90, C["white"], C["line"], 0.8)
    add_text(slide, 1.18, 2.06, 2.80, 0.20, "造型与内饰", 16, color=C["charcoal"], bold=True)
    add_bullets(
        slide,
        1.18,
        2.46,
        4.80,
        [
            "弱化商务车感，增强更年轻、更清晰的设计表达",
            "内饰同时满足简洁、高级感、科技感与长时间乘坐舒适",
            "中控、屏幕和整体座舱风格要更统一，减少“单调”和“粗糙感”",
            "具体前脸、侧面、尾部部位反馈，需问卷进一步验证",
        ],
        12.5,
        C["charcoal"],
        C["orange"],
        0.50,
    )
    add_rect(slide, 6.62, 1.80, 5.58, 4.90, C["soft_gray"], C["line"], 0.8)
    add_text(slide, 6.84, 2.06, 2.80, 0.20, "配置与体验", 16, color=C["charcoal"], bold=True)
    add_table(
        slide,
        6.84,
        2.42,
        4.92,
        3.62,
        ["优先强化", "说明"],
        [
            ["真实续航感知", "回答用户的补能安心与日常使用信心"],
            ["动力与加速体验", "体现新能源应有的爽感差异"],
            ["座椅舒适与空间灵活性", "覆盖家庭与长途场景"],
            ["车机与 OTA 能力", "把智能化从概念变成可感知体验"],
        ],
        [1.5, 2.3],
        10.5,
        C["blue"],
    )
    add_text(slide, 6.86, 6.18, 4.90, 0.22, "6/7 座接受度和各功能优先级仍需问卷验证。", 10.5, color=C["gray"])

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 13, "销售与数字化建议：体验竞争发生在看车之前，也发生在进店之后", "SERVICE & DIGITAL")
    metric_card(slide, 0.98, 1.82, 3.44, 1.38, "销售沟通", "减少只背参数的表达，增加顾客场景代入、真实对比解释与试驾重点引导。", C["orange"], C["soft_orange"])
    metric_card(slide, 4.72, 1.82, 3.44, 1.38, "试驾重点", "重点展示动力响应、操控和转向、舒适性、空间与座椅感受。", C["teal"], C["soft_teal"])
    metric_card(slide, 8.46, 1.82, 3.44, 1.38, "数字化承接", "官网、车型官网和 APP 应提前进入用户了解路径，承接预约、亮点和远程控车理解。", C["blue"], C["soft_blue"])
    add_table(
        slide,
        0.98,
        3.64,
        10.96,
        2.80,
        ["触点", "用户当前感受", "建议动作"],
        [
            ["门店讲解", "只讲参数，不够有温度", "统一场景化话术与顾客视角讲解方式"],
            ["试驾路线", "体验亮点没有被明确放大", "把动力、静音、空间和车机操作嵌入固定试驾脚本"],
            ["官网 / APP", "多数用户并不了解完整能力", "前置预约试驾、亮点说明、品牌文化和远程控车价值"],
        ],
        [1.0, 2.2, 3.0],
        11,
        C["charcoal"],
    )

    slide = prs.slides.add_slide(blank)
    page_chrome(slide, 14, "行动路径：先补齐基础体验，再进入问卷验证和更细的版本策略", "ROADMAP")
    stages = [
        ("阶段 1", "补齐基础体验", "续航可信、动力爽感、座舱科技感和服务温度作为第一优先级。"),
        ("阶段 2", "统一前端表达", "围绕家庭出游 + 城市通勤 + 长途安心，明确“适合谁、解决什么、为什么值得买”。"),
        ("阶段 3", "进入问卷验证", "验证 6/7 座、OTA/控车/泊车、价格与金融政策等具体优先级。"),
    ]
    y = 2.02
    for idx, (stage, title, note) in enumerate(stages, start=1):
        add_rect(slide, 1.08, y, 0.76, 0.52, C["orange"] if idx != 2 else C["teal"], C["orange"] if idx != 2 else C["teal"], 0)
        add_text(slide, 1.20, y + 0.12, 0.52, 0.16, f"{idx}", 12, color=C["white"], bold=True, align=PP_ALIGN.CENTER, font_name=FONT_NUM)
        add_rect(slide, 2.04, y - 0.02, 9.40, 0.94, C["white"], C["line"], 0.8)
        add_text(slide, 2.26, y + 0.10, 1.60, 0.18, stage, 13, color=C["gray"], bold=True)
        add_text(slide, 3.40, y + 0.08, 2.10, 0.20, title, 16, color=C["charcoal"], bold=True)
        add_text(slide, 5.62, y + 0.08, 5.44, 0.34, note, 12.5, color=C["charcoal"])
        if idx < len(stages):
            add_line(slide, 1.46, y + 0.92, 0.03, C["line"], 0.34)
        y += 1.38
    add_text(slide, 1.08, 6.20, 10.80, 0.18, "提醒：具体配置梯度、价格和政策影响占比，不在本次座谈中直接定量下结论。", 10.5, color=C["gray"])

    slide = prs.slides.add_slide(blank)
    set_bg(slide, C["dark"])
    add_text(slide, 1.02, 1.26, 2.60, 0.20, "CLOSING SUMMARY", 12, color=RGBColor(0xCA, 0xD5, 0xDA), bold=True, font_name=FONT_NUM)
    add_text(
        slide,
        1.02,
        2.08,
        10.80,
        1.20,
        "如果产品想真正打动这类用户，\n就必须同时回答“更安全、更智能、更好开，也更适合真实生活场景”这四个问题。",
        27,
        color=C["white"],
        bold=True,
    )
    add_text(
        slide,
        1.04,
        4.42,
        8.96,
        0.54,
        "本版为原生可编辑对象：标题、表格、流程、矩阵、建议模块均可直接在 PowerPoint 中二次修改。后续可以继续按模板缩略图细修单页版面。",
        14,
        color=RGBColor(0xD8, 0xE2, 0xE6),
    )
    add_line(slide, 0.88, 6.72, 11.52, RGBColor(0x77, 0x83, 0x86), 0.01)
    add_text(slide, 0.90, 6.76, 9.0, 0.14, SOURCE, 9, color=RGBColor(0xD8, 0xE2, 0xE6))
    add_text(slide, 12.00, 6.76, 0.30, 0.14, "15", 9, color=C["white"], align=PP_ALIGN.RIGHT, font_name=FONT_NUM)

    OUT.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(OUT))


if __name__ == "__main__":
    build()
    export_preview(OUT)
    print(OUT)
